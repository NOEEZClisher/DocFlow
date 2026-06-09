from __future__ import annotations

from pathlib import Path
from typing import Any

from docflow.core.document_model import Document, DocumentLoadError
from docflow.core.markdown_renderer import render_markdown

from .base import DocumentConverter


class PptxConverter(DocumentConverter):
    """Convert PPTX slide text into Markdown sections."""

    supported_extensions = frozenset({".pptx"})

    def convert(self, path: Path) -> Document:
        try:
            from pptx import Presentation
        except ModuleNotFoundError as exc:
            raise DocumentLoadError("python-pptx 패키지가 설치되어 있지 않습니다.") from exc

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise DocumentLoadError(f"PPTX 파일을 읽을 수 없습니다: {path}") from exc

        raw_text = "\n\n".join(
            _slide_to_markdown(slide, slide_number)
            for slide_number, slide in enumerate(presentation.slides, start=1)
        ).strip()
        return Document(
            path=path,
            kind="pptx",
            raw_text=raw_text,
            rendered_html=render_markdown(raw_text),
        )


def _slide_to_markdown(slide: Any, slide_number: int) -> str:
    texts = _slide_texts(slide)
    body = "\n\n".join(texts) if texts else "텍스트 없음"
    return f"## 슬라이드 {slide_number}\n\n{body}"


def _slide_texts(slide: Any) -> list[str]:
    texts: list[str] = []
    seen: set[str] = set()

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        for paragraph in shape.text_frame.paragraphs:
            text = _sanitize_markdown_text(paragraph.text)
            if not text:
                continue

            normalized = _normalize_for_deduplication(text)
            if normalized in seen:
                continue

            seen.add(normalized)
            texts.append(text)

    return texts


def _sanitize_markdown_text(text: str) -> str:
    text = " ".join(part.strip() for part in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"))
    return _escape_markdown_text(text.strip())


def _normalize_for_deduplication(text: str) -> str:
    return " ".join(text.split()).casefold()


def _escape_markdown_text(text: str) -> str:
    markdown_chars = "\\`*_{}[]()#+-.!|"
    escaped = []
    for char in text:
        if char in markdown_chars:
            escaped.append(f"\\{char}")
        else:
            escaped.append(char)
    return "".join(escaped)
