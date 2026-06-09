from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from docflow.converters.pptx_converter import PptxConverter
from docflow.converters.registry import default_registry


def save_presentation(path: Path, presentation: Any) -> Path:
    presentation.save(path)
    return path


def add_textbox(slide: Any, text: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text_frame.text = text


def test_pptx_converter_converts_single_slide_text(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_textbox(slide, "제출 안내")
    add_textbox(slide, "본문 텍스트")
    path = save_presentation(tmp_path / "single.pptx", presentation)

    document = PptxConverter().convert(path)

    assert document.kind == "pptx"
    assert document.raw_text == "## 슬라이드 1\n\n제출 안내\n\n본문 텍스트"
    assert document.rendered_html is not None
    assert "<h2>슬라이드 1</h2>" in document.rendered_html


def test_pptx_converter_converts_multiple_slides(tmp_path: Path) -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_textbox(first, "첫 번째 슬라이드")
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_textbox(second, "두 번째 슬라이드")
    path = save_presentation(tmp_path / "multi.pptx", presentation)

    document = PptxConverter().convert(path)

    assert "## 슬라이드 1" in document.raw_text
    assert "첫 번째 슬라이드" in document.raw_text
    assert "## 슬라이드 2" in document.raw_text
    assert "두 번째 슬라이드" in document.raw_text


def test_pptx_converter_marks_empty_slides(tmp_path: Path) -> None:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    path = save_presentation(tmp_path / "empty.pptx", presentation)

    document = PptxConverter().convert(path)

    assert document.raw_text == "## 슬라이드 1\n\n텍스트 없음"


def test_pptx_converter_escapes_markdown_text_and_skips_duplicates(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_textbox(slide, "중복 *강조*")
    add_textbox(slide, "중복 *강조*")
    path = save_presentation(tmp_path / "duplicates.pptx", presentation)

    document = PptxConverter().convert(path)

    assert document.raw_text == "## 슬라이드 1\n\n중복 \\*강조\\*"


def test_pptx_converter_sanitizes_line_breaks_and_pipe_characters(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_textbox(slide, "첫 줄\n둘째 줄 | 확인")
    path = save_presentation(tmp_path / "linebreaks.pptx", presentation)

    document = PptxConverter().convert(path)

    assert document.raw_text == "## 슬라이드 1\n\n첫 줄\n\n둘째 줄 \\| 확인"


def test_pptx_converter_keeps_common_sentence_punctuation(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_textbox(slide, "마감일: 6.10! 다시 제출-확인.")
    path = save_presentation(tmp_path / "punctuation.pptx", presentation)

    document = PptxConverter().convert(path)

    assert document.raw_text == "## 슬라이드 1\n\n마감일: 6.10! 다시 제출-확인."


def test_registry_selects_pptx_converter_for_pptx_files() -> None:
    converter = default_registry.get_converter(Path("slides.pptx"))

    assert isinstance(converter, PptxConverter)
